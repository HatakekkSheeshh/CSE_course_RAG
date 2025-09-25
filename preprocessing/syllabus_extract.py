# syllabus_extract.py
import fitz, re
from pptx import Presentation

def _kv(text, keys):
    for k in keys:
        m = re.search(rf'{k}\s*[:\-]\s*(.+)', text, re.I)
        if m: return m.group(1).strip()
    return None

def from_pdf(path):
    doc = fitz.open(path)
    text = "\n".join(page.get_text() for page in doc)
    return {
        "course_name": _kv(text, ["Course Name","Course Title","Title"]),
        "instructor":  _kv(text, ["Instructor","Lecturer","Teacher"]),
        "etcs":        _kv(text, ["ETCS","Credits"]),
        "semester":    _kv(text, ["Semester","Term"]),
    }

def from_pptx(path):
    prs = Presentation(path)
    txt = []
    for s in prs.slides:
        if s.shapes.title: txt.append(s.shapes.title.text)
        for sh in s.shapes:
            if getattr(sh, "has_text_frame", False):
                txt.append(sh.text)
        if s.has_notes_slide and s.notes_slide.notes_text_frame:
            txt.append(s.notes_slide.notes_text_frame.text)
    text = "\n".join(txt)
    return {
        "course_name": _kv(text, ["Course Name","Course Title","Title"]),
        "instructor":  _kv(text, ["Instructor","Lecturer","Teacher"]),
        "ects":        _kv(text, ["ECTS","Credits"]),
        "semester":    _kv(text, ["Semester","Term"]),
    }
